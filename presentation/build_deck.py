# -*- coding: utf-8 -*-
"""Генератор презентации Halyk Smart Travel Companion (.pptx).

Запуск:  .venv\\Scripts\\python.exe presentation\\build_deck.py
Выход:   presentation/Halyk_Smart_Travel.pptx  (16:9, 21 слайд)

Дек пересобирается из этого скрипта — правьте текст здесь и запускайте заново.
Фирменные цвета взяты из frontend/src/styles/companion.css.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ── Палитра Halyk ──
# ── тёплая палитра в эстетике Anthropic: кремовый холст, тёплый тёмный текст, зелёный Halyk ──
GREEN   = RGBColor(0x00, 0xA8, 0x4B)   # бренд-зелёный (заливки, полосы)
GREEN_D = RGBColor(0x05, 0x7A, 0x39)   # приглушённый зелёный (акцентный текст, рамки)
GREEN_L = RGBColor(0xE7, 0xF1, 0xE7)   # светло-зелёная плашка (тёплая)
GREEN_X = RGBColor(0xF1, 0xF6, 0xEF)   # едва-зелёная плашка
INK     = RGBColor(0x23, 0x21, 0x1C)   # тёплый почти-чёрный (заголовки/текст)
DARK    = RGBColor(0x2A, 0x27, 0x20)   # тёплый тёмный (только для мелких акцентов)
T2      = RGBColor(0x5C, 0x57, 0x4E)   # тёплый серый (вторичный текст)
T3      = RGBColor(0x99, 0x93, 0x86)   # тёплый светло-серый (подписи)
BG      = RGBColor(0xF0, 0xED, 0xE4)
BORDER  = RGBColor(0xE6, 0xE0, 0xD2)   # тёплая рамка
HAIR    = RGBColor(0xED, 0xE8, 0xDC)   # волосяная линия внутри карточек
AMBER   = RGBColor(0xC8, 0x7E, 0x3C)   # тёплый янтарь (под Anthropic clay)
AMBER_BG= RGBColor(0xF7, 0xEF, 0xE2)
RED     = RGBColor(0xC0, 0x53, 0x43)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xFD, 0xFC, 0xF9)   # тёплый белый — карточки на кремовом холсте
PAPER   = RGBColor(0xF4, 0xF1, 0xE9)   # холст контентных слайдов (тёплый крем)
HERO    = RGBColor(0xEF, 0xEB, 0xDF)   # чуть глубже крем — титул/финал
ROW     = RGBColor(0xF5, 0xF2, 0xEA)   # чётные строки таблицы

FONT  = "Segoe UI"      # основной текст
SERIF = "Georgia"       # заголовки в эстетике Anthropic

EMU_IN = 914400
PW, PH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(PW)
prs.slide_height = Inches(PH)
BLANK = prs.slide_layouts[6]


# ── базовые помощники ──
def slide():
    return prs.slides.add_slide(BLANK)


def set_bg(s, color):
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = color


def _no_line(sh):
    sh.line.fill.background()


def soft_shadow(sh, blur=0.14, dist=0.05, alpha=88, color="2B2620"):
    """Мягкая тёплая тень для карточек/изображений (alpha — прозрачность %, больше = бледнее)."""
    from pptx.oxml.ns import qn
    spPr = sh._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    shdw = eff.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur * EMU_IN)),
        'dist': str(int(dist * EMU_IN)),
        'dir': '5400000',          # вниз
        'rotWithShape': '0',
    })
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': color})
    al = clr.makeelement(qn('a:alpha'), {'val': str(int((100 - alpha) * 1000))})
    clr.append(al); shdw.append(clr); eff.append(shdw); spPr.append(eff)
    return sh


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08, shadow=False):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shape)
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    if shadow:
        soft_shadow(shape)
    return shape


def card(s, l, t, w, h, fill=CARD_BG, line=BORDER, radius=0.06, shadow=True):
    """Стандартная белая карточка: тонкая рамка + мягкая тень → «лежит» на холсте."""
    return rect(s, l, t, w, h, fill=fill, line=line, line_w=1.0, rounded=True, radius=radius, shadow=shadow)


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, runs, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, first=False, line=1.05, font=FONT):
    """runs: str | list[(text,color,bold) | (text,color) | text]."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for item in runs:
        if isinstance(item, str):
            txt, c, b = item, color, bold
        elif len(item) == 2:
            txt, c, b = item[0], item[1], bold
        else:
            txt, c, b = item
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.name = font
        r.font.color.rgb = c
    return p


# ── фирменные элементы ──
def wordmark(s, l, t, scale=1.0, on_dark=False):
    """Зелёный квадрат с «H» + текст «Halyk»."""
    sq = scale * 0.42
    box = rect(s, l, t, sq, sq, fill=GREEN, rounded=True, radius=0.30)
    box.shadow.inherit = False
    _, tf = textbox(s, l, t, sq, sq, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("H", WHITE, True)], size=int(20 * scale), align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    _, tf2 = textbox(s, l + sq + 0.12 * scale, t - 0.04 * scale, 3.2 * scale, sq + 0.1,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf2, [("Halyk ", WHITE if on_dark else INK, True),
               ("Smart Travel", GREEN, True)],
         size=int(17 * scale), first=True, space_after=0)


def top_accent(s):
    rect(s, 0, 0, PW, 0.05, fill=GREEN)   # тонкая брендовая нить


def footer(s, n):
    _, tf = textbox(s, PW - 3.2, PH - 0.45, 3.0, 0.3, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("Halyk Smart Travel Companion   ·   ", T3),
              (f"{n:02d}", T2)], size=9, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def kicker(s, text, x=0.9, y=0.66, color=GREEN_D):
    _, tf = textbox(s, x, y, 11, 0.4)
    para(tf, [(text.upper(), color, True)], size=11.5, first=True, space_after=0)


def title(s, text, x=0.9, y=1.04, w=11.5, size=33, color=INK):
    _, tf = textbox(s, x, y, w, 1.2)
    para(tf, [(text, color, False)], size=size, first=True, space_after=0, line=1.02, font=SERIF)


def title_rule(s, x=0.9, y=1.78, w=0.9):
    rect(s, x, y, w, 0.04, fill=GREEN)


def bullets(s, l, t, w, h, items, size=16, gap=11, marker="•", mcolor=GREEN, anchor=MSO_ANCHOR.TOP):
    _, tf = textbox(s, l, t, w, h, anchor=anchor)
    for i, it in enumerate(items):
        # it: str | (head, tail) where head bold ink, tail T2
        if isinstance(it, tuple):
            runs = [(f"{marker}  ", mcolor, True), (it[0], INK, True)]
            if it[1]:
                runs.append((it[1], T2, False))
        else:
            runs = [(f"{marker}  ", mcolor, True), (it, INK, False)]
        para(tf, runs, size=size, first=(i == 0), space_after=gap, line=1.08)
    return tf


def place_image(s, path, box_l, box_t, box_w, box_h, frame=True, align="center", valign="middle"):
    """Вписать изображение в рамку box, сохранив пропорции. Возвращает (l,t,w,h) дюймов."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w; h = box_w / ar
    else:
        h = box_h; w = box_h * ar
    if align == "center":
        l = box_l + (box_w - w) / 2
    elif align == "left":
        l = box_l
    else:
        l = box_l + (box_w - w)
    if valign == "middle":
        t = box_t + (box_h - h) / 2
    elif valign == "top":
        t = box_t
    else:
        t = box_t + (box_h - h)
    if frame:
        fr = rect(s, l - 0.06, t - 0.06, w + 0.12, h + 0.12, fill=WHITE,
                  line=BORDER, line_w=1.0, rounded=True, radius=0.03)
        soft_shadow(fr, blur=0.12, dist=0.07, alpha=82)
    pic = s.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(0.75)
    return l, t, w, h


def chip(s, l, t, text, fill=GREEN_L, color=GREEN_D, w=None, size=12):
    w = w or (0.18 + 0.105 * len(text))
    c = rect(s, l, t, w, 0.38, fill=fill, rounded=True, radius=0.5)
    c.shadow.inherit = False
    _, tf = textbox(s, l, t, w, 0.38, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(text, color, True)], size=size, align=PP_ALIGN.CENTER, first=True, space_after=0)
    return w


def asset(name):
    return os.path.join(ASSETS, name)


# ══════════════════════════════════════════════════════════════════════
# СЛАЙДЫ
# ══════════════════════════════════════════════════════════════════════

# 1 — Титул
def s01():
    s = slide(); set_bg(s, HERO)
    rect(s, 0, 0, PW, 0.05, fill=GREEN)
    wordmark(s, 0.9, 0.72, scale=1.15)
    _, tf = textbox(s, 0.9, 2.45, 11.6, 2.6)
    para(tf, [("Smart Travel ", INK), ("Companion", GREEN_D)],
         size=54, first=True, space_after=12, line=1.0, font=SERIF)
    para(tf, "AI-ассистент, который превращает покупку билета", size=21, color=T2, space_after=2)
    para(tf, "в операционный штаб семейной поездки", size=21, color=T2, space_after=0)
    # нижняя плашка
    chip(s, 0.9, 5.6, "Алматы → Астана", size=13)
    chip(s, 3.3, 5.6, "ночной поезд · семья из 4", size=13)
    _, tf2 = textbox(s, 0.9, 6.72, 11.5, 0.5)
    para(tf2, [("Демо-прототип  ·  LangGraph + claude-haiku-4-5  ·  один контейнер, end-to-end", T3)],
         size=13, first=True, space_after=0)


# 2 — Проблема
def s02():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 2)
    kicker(s, "01 · Проблема")
    title(s, "Клиент купил билет — и ушёл"); title_rule(s)
    bullets(s, 0.9, 2.1, 6.6, 4.5, [
        ("Halyk Travel сегодня — ", "сервис покупки билетов. На этом контакт с клиентом заканчивается."),
        ("Отель, страховка, трансфер, развлечения, питание ", "клиент организует сам, в чужих приложениях."),
        ("Все downstream-транзакции ", "уходят к конкурентам — вместе с данными и лояльностью."),
    ], size=17, gap=16)
    # правый акцент — белая карточка с тенью и зелёной полосой слева
    card(s, 8.0, 2.15, 4.43, 4.05, fill=CARD_BG, radius=0.05)
    rect(s, 8.0, 2.32, 0.10, 3.71, fill=GREEN)
    _, tf = textbox(s, 8.45, 2.6, 3.6, 3.2, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("1", GREEN_D)], size=80, first=True, space_after=2, line=1.0, font=SERIF)
    para(tf, [("транзакция вместо ", INK, True), ("15+", GREEN, True)], size=19, space_after=12)
    para(tf, "Каждая неоформленная услуга — упущенная выручка и ушедший клиент.",
         size=13.5, color=T2, line=1.15)


# 3 — Инсайт
def s03():
    s = slide(); set_bg(s, HERO); top_accent(s); footer(s, 3)
    kicker(s, "02 · Инсайт")
    _, tf = textbox(s, 0.9, 1.7, 11.5, 2.6)
    para(tf, [("Продукт строится на ", INK), ("карте тревог клиента,", GREEN_D)],
         size=40, first=True, space_after=4, line=1.08, font=SERIF)
    para(tf, [("а не на списке сервисов", INK)], size=40, space_after=0, font=SERIF)
    bullets(s, 0.9, 4.5, 11.5, 2.4, [
        "Каждое сообщение отвечает на конкретную мысль клиента в конкретный момент.",
        "«Когда предлагать Appteka» и «когда родитель думает об аптечке» — разные вопросы.",
        "Клиент не чувствует, что ему продают, — он чувствует, что о нём заботятся.",
    ], size=17, gap=15)


# 4 — Решение
def s04():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 4)
    kicker(s, "03 · Решение")
    title(s, "Из разовой покупки — в штаб всей поездки"); title_rule(s)
    bullets(s, 0.9, 2.05, 11.5, 1.5, [
        ("Halyk ведёт клиента ", "от покупки билета до возвращения домой, закрывая каждую тревогу в нужный момент."),
        ("~18 уведомлений за 18 дней ", "— примерно одно в день. Умное молчание важнее умных сообщений."),
    ], size=16.5, gap=12)
    # лента пути
    steps = ["🎫 Билет", "🏨 Жильё", "📋 Документы", "🛡️ Страховка", "💰 Бюджет",
             "💊 Аптечка", "🚖 Трансфер", "🎬 Досуг", "🏁 Итоги"]
    n = len(steps); x = 0.9; y = 4.4; gap = 0.16
    cw = (PW - 1.8 - gap * (n - 1)) / n
    for i, st in enumerate(steps):
        first = (i == 0); last = (i == n - 1)
        fill = GREEN if first else (GREEN_L if not last else GREEN_D)
        tcol = WHITE if (first or last) else GREEN_D
        c = rect(s, x, y, cw, 1.0, fill=fill, rounded=True, radius=0.14)
        c.shadow.inherit = False
        _, tf = textbox(s, x, y, cw, 1.0, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(st, tcol, True)], size=10.5, align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.0)
        if not last:
            _, ta = textbox(s, x + cw - 0.02, y, gap + 0.04, 1.0, anchor=MSO_ANCHOR.MIDDLE)
            para(ta, [("→", T3, True)], size=14, align=PP_ALIGN.CENTER, first=True, space_after=0)
        x += cw + gap
    _, tl = textbox(s, 0.9, 5.7, 11.5, 0.4)
    para(tl, [("Фаза 0 (Т−14)   →   Фаза 1 (Т−7…−3)   →   Фаза 2 (в поезде)   →   Фаза 3 (в Астане)   →   Фаза 4 (Итоги)", T2)],
         size=13, align=PP_ALIGN.CENTER, first=True, space_after=0)


# 5 — Референсный сценарий
def s05():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 5)
    kicker(s, "Референсный сценарий")
    title(s, "Одна семья, выходные в Астане"); title_rule(s)
    # карточка семьи
    card(s, 0.9, 2.1, 6.5, 4.5, radius=0.05)
    people = [("А", "Айдар", "организатор, плательщик"),
              ("Ал", "Алия", "удостоверение истекает 28 июля"),
              ("9", "Айша", "9 лет"),
              ("5", "Тимур", "5 лет")]
    y = 2.45
    for av, nm, sub in people:
        a = rect(s, 1.2, y, 0.5, 0.5, fill=GREEN_L, rounded=True, radius=0.5)
        a.shadow.inherit = False
        _, tfa = textbox(s, 1.2, y, 0.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        para(tfa, [(av, GREEN_D, True)], size=13, align=PP_ALIGN.CENTER, first=True, space_after=0)
        _, tfn = textbox(s, 1.85, y - 0.04, 5.2, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, [(nm + "   ", INK, True), (sub, T2, False)], size=14.5, first=True, space_after=0)
        y += 0.72
    _, tfr = textbox(s, 1.2, y + 0.05, 5.9, 0.9)
    para(tfr, [("Маршрут:  ", T2), ("Алматы → Астана, ночной поезд ~13 ч", INK, True)], size=13.5, first=True, space_after=5)
    para(tfr, [("Даты:  ", T2), ("4→5 июня туда · 5–7 июня в Астане · дождливая суббота +14°", INK, True)], size=13.5, space_after=0)
    # правая колонка
    bullets(s, 7.8, 2.3, 4.6, 2.0, [
        ("Главный страх организатора: ", "«что-то не предусмотреть»."),
        ("Готов платить ", "за снятие тревоги, а не за сами сервисы."),
    ], size=16, gap=14)
    place_image(s, asset("w1-wizard-step1.png"), 7.8, 4.2, 4.6, 2.3, align="center", valign="top")
    _, tc = textbox(s, 7.8, 6.5, 4.6, 0.3)
    para(tc, [("Демо открывается с покупки билета (3 шага)", T3)], size=10.5, align=PP_ALIGN.CENTER, first=True, space_after=0)


# 6 — Карта тревог (таблица)
def s06():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 6)
    kicker(s, "04 · Что видит клиент")
    title(s, "Карта тревог клиента", size=30); title_rule(s)
    rows = [
        ("Момент", "Мысль в голове", "Решение Halyk"),
        ("Сразу после билета", "«Надо отель — разберут быстро»", "Booking.com"),
        ("После отеля", "«Документы в порядке?»", "66 госуслуг"),
        ("После отеля", "«Страховка нужна, наверное»", "СК Халык"),
        ("После отеля", "«Сколько вообще уйдёт денег?»", "Автосмета"),
        ("За 7 дней", "«Аптечку собрать надо»", "Appteka"),
        ("За 3 дня", "«Как с вокзала доедем?»", "inDrive"),
        ("За 3 дня", "«Чем занять детей в дождь?»", "Kino.kz"),
        ("В поездке", "«Сколько уже потратили?»", "Трекер бюджета"),
        ("После поездки", "«Сколько вышло в итоге?»", "Итоги план/факт"),
    ]
    L, T, W = 0.9, 2.05, 11.5
    rowh = 0.475
    widths = [3.0, 5.4, 3.1]
    # таблица-карточка: тёплый белый фон + мягкая тень
    rect(s, L, T, W, rowh * len(rows), fill=CARD_BG, line=BORDER, line_w=1.0, rounded=True, radius=0.02, shadow=True)
    for r, row in enumerate(rows):
        y = T + r * rowh
        head = (r == 0)
        if head:
            rect(s, L, y, W, rowh, fill=GREEN, line=None)
        elif r % 2 == 0:
            rect(s, L + 0.02, y, W - 0.04, rowh, fill=ROW, line=None)
        x = L
        for c, cell in enumerate(row):
            _, tf = textbox(s, x + 0.22, y, widths[c] - 0.34, rowh, anchor=MSO_ANCHOR.MIDDLE)
            if head:
                para(tf, [(cell, WHITE, True)], size=13, first=True, space_after=0)
            elif c == 2:
                para(tf, [(cell, GREEN_D, True)], size=13.5, first=True, space_after=0)
            else:
                para(tf, [(cell, INK if c == 0 else T2, c == 0)], size=13, first=True, space_after=0)
            x += widths[c]


# 7 — Путь 10/5
def s07():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 7)
    kicker(s, "04 · Что видит клиент")
    title(s, "Путь клиента: 10 этапов, 5 фаз"); title_rule(s)
    phases = [
        ("Фаза 0", "Т−14", ["Жильё", "Документы", "Страховка", "Бюджет"]),
        ("Фаза 1", "Т−7…−3", ["Аптечка", "Трансфер", "Развлечения", "Ресторан", "Погода"]),
        ("Фаза 2", "В поезде", ["Такси с вокзала", "Продукты (апарт)"]),
        ("Фаза 3", "В Астане", ["Трекер бюджета", "Напоминания", "Экстренное", "Сувениры"]),
        ("Фаза 4", "Итоги", ["План / Факт", "Бонусы Halyk+", "Flywheel"]),
    ]
    n = len(phases); x = 0.9; y = 2.3; gap = 0.2
    cw = (PW - 1.8 - gap * (n - 1)) / n
    ch = 3.9
    for i, (ph, when, items) in enumerate(phases):
        # единая белая карточка с тенью, сверху — цветная «шапка»
        card(s, x, y, cw, ch, radius=0.07)
        head = rect(s, x, y, cw, 0.86, fill=GREEN if i == 0 else GREEN_D, rounded=True, radius=0.08)
        head.shadow.inherit = False
        rect(s, x, y + 0.5, cw, 0.36, fill=GREEN if i == 0 else GREEN_D)   # выпрямляем низ шапки
        _, tfh = textbox(s, x, y + 0.08, cw, 0.74, anchor=MSO_ANCHOR.MIDDLE)
        para(tfh, [(ph, WHITE, True)], size=15, align=PP_ALIGN.CENTER, first=True, space_after=1)
        para(tfh, [(when, RGBColor(0xD1,0xFA,0xE0), False)], size=11, align=PP_ALIGN.CENTER, space_after=0)
        _, tfb = textbox(s, x + 0.18, y + 1.12, cw - 0.36, ch - 1.25)
        for j, it in enumerate(items):
            para(tfb, [("• ", GREEN, True), (it, INK, False)], size=12.5, first=(j == 0), space_after=9, line=1.05)
        if i < n - 1:
            _, ta = textbox(s, x + cw - 0.04, y, gap + 0.08, 0.86, anchor=MSO_ANCHOR.MIDDLE)
            para(ta, [("→", T3, True)], size=15, align=PP_ALIGN.CENTER, first=True, space_after=0)
        x += cw + gap
    _, tf = textbox(s, 0.9, 6.45, 11.5, 0.4)
    para(tf, [("Каждый пункт: ", T2), ("locked → wait → done", GREEN_D, True),
              (". Разблокировка — автоматически при смене фазы.", T2)],
         size=12.5, align=PP_ALIGN.CENTER, first=True, space_after=0)


# ── шаблон слайда «скриншот + тезисы» ──
def shot_slide(n, kick, ttl, img, bul, caption="", img_w=5.6, star=False):
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, n)
    kicker(s, kick, color=AMBER if star else GREEN)
    title(s, ttl); title_rule(s, w=1.1)
    place_image(s, asset(img), 0.9, 2.15, img_w, 4.7, align="left", valign="top")
    if caption:
        _, tc = textbox(s, 0.9, 2.15 + 4.7 + 0.02, img_w, 0.3)
        para(tc, [(caption, T3)], size=10.5, align=PP_ALIGN.CENTER, first=True, space_after=0)
    bx = 0.9 + img_w + 0.65
    bullets(s, bx, 2.15, PW - bx - 0.7, 4.7, bul, size=16.5, gap=16, anchor=MSO_ANCHOR.MIDDLE)
    return s


# 8 — Вход (демо)
def s08():
    shot_slide(8, "04 · Что видит клиент", "Билет — точка входа. Дальше AI собирает поездку",
               "w2-companion-start.png", [
        ("3-шаговый визард ", "покупки билета → передаёт клиента в companion."),
        ("Первое сообщение: ", "билеты куплены, факты поездки, «я — ваш ассистент»."),
        ("Второе: ", "вопрос про жильё — Booking.com и бонусы Halyk."),
        ("Справа — живой Travel Plan ", "из 10 пунктов и бюджет."),
    ], caption="Старт companion после покупки билета")


# 9 — УТП документы ★
def s09():
    s = shot_slide(9, "★ Главное конкурентное преимущество",
               "Проактивная проверка документов",
               "w3-doc-alert-insurance.png", [
        ("Halyk сам проверяет ", "документы семьи через 66 госуслуг — молча."),
        ("Замечает: ", "удостоверение Алии истекает 28 июля → предлагает eGov заранее."),
        ("Только Halyk знает эти данные. ", "Конкурент этого не сделает."),
        ("Одно такое уведомление ", "строит больше доверия, чем любой кешбэк."),
    ], caption="Алерт о документах появляется сам", star=True)


# 10 — Бюджет
def s10():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 10)
    kicker(s, "04 · Что видит клиент")
    title(s, "Сквозной бюджет: «всё под контролем»"); title_rule(s)
    # крупная панель бюджета
    place_image(s, asset("budget_panel.png"), 0.9, 2.5, 5.4, 2.2, align="left", valign="top")
    _, tc = textbox(s, 0.9, 4.85, 5.4, 0.3)
    para(tc, [("Панель бюджета на старте поездки", T3)], size=10.5, align=PP_ALIGN.CENTER, first=True, space_after=0)
    # трасса роста факта
    rect(s, 0.9, 5.4, 5.4, 1.0, fill=GREEN_X, line=GREEN_L, rounded=True, radius=0.08, shadow=True)
    _, tt = textbox(s, 1.2, 5.5, 5.0, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tt, [("Факт растёт:  ", T2), ("38 000 → 89 000 → 105 000 ₸", GREEN_D, True)], size=14, first=True, space_after=3)
    para(tt, [("Расчётное тает, итог держится ", T2), ("~175 000 ₸", INK, True)], size=14, space_after=0)
    bullets(s, 6.9, 2.5, 5.6, 4.0, [
        ("Два уровня: ", "Факт (оплачено/зарезервировано) + Расчётное (прогноз системы)."),
        ("Смета собирается сама ", "по известным данным — клиент ничего не вводит."),
        ("По мере оплат ", "деньги перетекают из расчётного в факт, итог стабилен."),
        ("Это эмоциональный продукт, ", "не аналитический: клиент перестаёт беспокоиться о деньгах."),
    ], size=16, gap=15)


# 11 — Погода
def s11():
    shot_slide(11, "04 · Что видит клиент", "Погода меняет план поездки",
               "03-phase1-done.png", [
        ("Система мониторит прогноз ", "с момента покупки билета."),
        ("Дождливая суббота +14° ", "→ только крытые активности (Думан), дождевики, такси."),
        ("Рекомендации адаптируются: ", "дождь · жара · холод · норма — разные сценарии."),
        ("В демо — дождевой сценарий: ", "самая богатая логика."),
    ], caption="Фаза 1: подготовка под дождь")


# 12 — Фазы 2-3
def s12():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 12)
    kicker(s, "04 · Что видит клиент")
    title(s, "В пути и на месте: момент важнее экрана"); title_rule(s)
    place_image(s, asset("04-phase3-done.png"), 0.9, 2.15, 5.5, 4.5, align="left", valign="top")
    _, tc = textbox(s, 0.9, 6.66, 5.5, 0.3)
    para(tc, [("Фаза 3: живой трекер и напоминания", T3)], size=10.5, align=PP_ALIGN.CENTER, first=True, space_after=0)
    bx = 7.1
    bullets(s, bx, 2.35, PW - bx - 0.7, 2.4, [
        ("За 10 минут до Астаны ", "— такси одним тапом (клиент стоит с детьми и чемоданами)."),
        ("Живой трекер: ", "платежи картой Halyk сами падают в категории бюджета."),
    ], size=15.5, gap=13)
    place_image(s, asset("emergency.png"), bx, 4.75, 5.2, 1.7, align="left", valign="top")
    _, te = textbox(s, bx, 4.45, 5.2, 0.3)
    para(te, [("Экстренный блок — всегда под рукой:", INK, True)], size=13.5, first=True, space_after=0)


# 13 — Итоги + Flywheel
def s13():
    shot_slide(13, "04 · Что видит клиент", "Итоги план/факт и новый виток",
               "05-results.png", [
        ("План 175 000 → Факт 169 500 ₸ ", "(−5 500). Уложились в бюджет 🎯."),
        ("Бонусы Halyk+ начислены, ", "категория клиента повышена."),
        ("Три варианта следующей поездки ", "— Бурабай, Алаколь, Шымкент."),
        ("Один тап → ", "Halyk Travel с поиском. Цикл (flywheel) замыкается."),
    ], caption="Экран Итоги + Flywheel", img_w=5.9)


# 14 — Ценность для Halyk
def s14():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 14)
    kicker(s, "05 · Ценность и будущее")
    title(s, "Каждая тревога — точка монетизации"); title_rule(s)
    items = [
        ("Booking.com", "revenue share + бонусы Halyk"),
        ("СК Халык", "страховая премия"),
        ("Appteka · inDrive · Airba", "revenue share с партнёров"),
        ("Kino.kz", "прямой доход (сервис Halyk) + бонусы"),
        ("Halyk Market", "продажа + рассрочка Homebank"),
        ("Halyk Pay / QR", "interchange + данные о тратах"),
        ("66 госуслуг", "доверие и retention"),
        ("Halyk+", "лояльность → flywheel"),
    ]
    L, T = 0.9, 2.05; gx, gy = 0.3, 0.24; cw, chh = (11.5 - gx) / 2, 0.82
    for i, (svc, mon) in enumerate(items):
        r, c = divmod(i, 2)
        x = L + c * (cw + gx); y = T + r * (chh + gy)
        card(s, x, y, cw, chh, radius=0.12)
        rect(s, x, y + 0.12, 0.10, chh - 0.24, fill=GREEN, rounded=True, radius=0.5)
        _, tf = textbox(s, x + 0.34, y, cw - 0.54, chh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(svc, INK, True)], size=15, first=True, space_after=2)
        para(tf, [(mon, T2, False)], size=12.5, space_after=0)
    _, tb = textbox(s, 0.9, 6.66, 11.5, 0.4)
    para(tb, [("Клиент остаётся внутри экосистемы Halyk всю поездку — не уходит к конкурентам.", GREEN_D, True)],
         size=13.5, align=PP_ALIGN.CENTER, first=True, space_after=0)


# 15 — Три вывода
def s15():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 15)
    kicker(s, "05 · Ценность и будущее")
    title(s, "Три ключевых вывода"); title_rule(s)
    cards = [
        ("1", "Тревоги, не сервисы", "Продукт отвечает на конкретную мысль клиента в нужный момент времени, а не предлагает услуги по календарю."),
        ("2", "Проактивная проверка", "Документы и погода — преимущество, доступное только Halyk: данные семьи уже в экосистеме."),
        ("3", "Трекер = стимул карты", "Полная картина трат только по транзакциям Halyk. Наличные не видны — клиент сам хочет платить картой."),
    ]
    x = 0.9; y = 2.3; gap = 0.4; cw = (11.5 - 2 * gap) / 3; chh = 3.9
    for num, head, body in cards:
        card(s, x, y, cw, chh, radius=0.05)
        circ = rect(s, x + 0.35, y + 0.45, 0.85, 0.85, fill=GREEN, rounded=True, radius=0.5)
        soft_shadow(circ, blur=0.06, dist=0.03, alpha=84)
        _, tn = textbox(s, x + 0.35, y + 0.45, 0.85, 0.85, anchor=MSO_ANCHOR.MIDDLE)
        para(tn, [(num, WHITE, True)], size=30, align=PP_ALIGN.CENTER, first=True, space_after=0)
        _, tf = textbox(s, x + 0.35, y + 1.5, cw - 0.7, chh - 1.7)
        para(tf, [(head, INK, True)], size=18, first=True, space_after=10, line=1.05)
        para(tf, [(body, T2, False)], size=14, space_after=0, line=1.12)
        x += cw + gap


# 16 — Масштабирование
def s16():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 16)
    kicker(s, "05 · Ценность и будущее")
    title(s, "От семьи на поезде — к B2B-командировкам"); title_rule(s)
    phases = [
        ("Фаза 1 · MVP", "Семья, ЖД, маршруты КЗ", "Жильё + страховка + трансфер + Kino.kz"),
        ("Фаза 2", "Авиа + международка", "Виза, медицина за рубежом, валюта в плане"),
        ("Фаза 3", "Персонализация и AI", "Рекомендации и бюджет на основе истории поездок"),
        ("Фаза 4", "B2B — командировки", "Корпоративный Travel Plan, отчёты, HR-интеграции"),
    ]
    x = 0.9; y = 2.3; gap = 0.25; cw = (11.5 - 3 * gap) / 4; chh = 3.8
    for i, (ph, head, body) in enumerate(phases):
        on = (i == 0)
        rect(s, x, y, cw, chh, fill=GREEN_L if on else CARD_BG, line=GREEN if on else BORDER,
             line_w=1.5 if on else 1.0, rounded=True, radius=0.06, shadow=True)
        _, th = textbox(s, x + 0.25, y + 0.35, cw - 0.5, 1.5)
        para(th, [(ph, GREEN_D if on else T2, True)], size=14, first=True, space_after=8)
        para(th, [(head, INK, True)], size=16.5, space_after=0, line=1.05)
        _, tb = textbox(s, x + 0.25, y + 1.9, cw - 0.5, chh - 2.1)
        para(tb, [(body, T2, False)], size=13, first=True, space_after=0, line=1.12)
        if i < 3:
            _, ta = textbox(s, x + cw - 0.02, y, gap + 0.04, chh, anchor=MSO_ANCHOR.MIDDLE)
            para(ta, [("→", T3, True)], size=16, align=PP_ALIGN.CENTER, first=True, space_after=0)
        x += cw + gap


# 17 — Живое демо
def s17():
    s = slide(); set_bg(s, HERO); top_accent(s); footer(s, 17)
    chip(s, 0.9, 0.8, "ДЕМО", size=12.5, w=1.05)
    _, tf = textbox(s, 0.9, 1.65, 7.2, 2.0)
    para(tf, [("Демо вживую", INK)], size=46, first=True, space_after=10, font=SERIF)
    para(tf, [("docker compose up  →  localhost:8000", GREEN_D, True)], size=18, space_after=0)
    bullets(s, 0.9, 3.7, 7.0, 2.8, [
        "Визард покупки билета → передача в companion.",
        "Жильё → проактивный алерт о документах → страховка → бюджет.",
        "Погодная адаптация, живой трекер, экстренный блок.",
        "Итоги план/факт и Flywheel.",
    ], size=16, gap=14)
    _, tn = textbox(s, 0.9, 6.4, 7.2, 0.5)
    para(tn, [("Работает полностью оффлайн — fallback на каждом шаге, сеть не нужна.", T3)],
         size=13, first=True, space_after=0)
    place_image(s, asset("02-phase0-done.png"), 8.5, 1.65, 4.0, 4.7, align="center", valign="top", frame=True)
    _, tc = textbox(s, 8.5, 6.45, 4.0, 0.3)
    para(tc, [("Резервный скриншот на случай сбоя", T3)], size=10, align=PP_ALIGN.CENTER, first=True, space_after=0)


# 18 — Архитектура
def s18():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 18)
    kicker(s, "06 · Как это построено")
    title(s, "Один контейнер, весь путь end-to-end"); title_rule(s)
    # диаграмма-конвейер
    boxes = [("React SPA", "Vite + TS"), ("Django + DRF", "5 endpoints"),
             ("LangGraph", "ядро-граф"), ("claude-haiku-4-5", "или fallback")]
    x = 0.9; y = 2.35; gap = 0.55; cw = (11.5 - 3 * gap) / 4; chh = 1.3
    for i, (h1, h2) in enumerate(boxes):
        fill = GREEN if i == 2 else CARD_BG
        tcol = WHITE if i == 2 else INK
        scol = RGBColor(0xD1,0xFA,0xE0) if i == 2 else T2
        rect(s, x, y, cw, chh, fill=fill, line=None if i == 2 else BORDER, rounded=True, radius=0.10, shadow=True)
        _, tf = textbox(s, x, y, cw, chh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(h1, tcol, True)], size=15.5, align=PP_ALIGN.CENTER, first=True, space_after=3)
        para(tf, [(h2, scol, False)], size=12, align=PP_ALIGN.CENTER, space_after=0)
        if i < 3:
            _, ta = textbox(s, x + cw, y, gap, chh, anchor=MSO_ANCHOR.MIDDLE)
            para(ta, [("→", GREEN, True)], size=20, align=PP_ALIGN.CENTER, first=True, space_after=0)
        x += cw + gap
    _, td = textbox(s, 0.9, 3.75, 11.5, 0.4)
    para(td, [("Django ORM (SQLite) — единый источник правды.  WhiteNoise отдаёт SPA.  Всё в одном Docker-образе.", T2)],
         size=13, align=PP_ALIGN.CENTER, first=True, space_after=0)
    bullets(s, 0.9, 4.6, 11.5, 2.0, [
        ("Стек: ", "Django 5 + DRF · React 18 / Vite / TS · SQLite · LangChain/LangGraph · Gunicorn · Docker."),
        ("Запуск: ", "docker compose up → localhost:8000, end-to-end в один шаг."),
        ("Демо работает полностью оффлайн: ", "у каждого AI-шага есть готовый fallback."),
    ], size=15.5, gap=13)


# 19 — LangGraph
def s19():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 19)
    kicker(s, "06 · Как это построено")
    title(s, "Граф — это и есть продукт"); title_rule(s)
    # граф узлов
    nodes = ["apply_answer", "advance", "generate"]
    x = 0.9; y = 2.4; gap = 0.7; cw = 2.7; chh = 1.0
    for i, nd in enumerate(nodes):
        rect(s, x, y, cw, chh, fill=GREEN_L, line=GREEN, rounded=True, radius=0.12, shadow=True)
        _, tf = textbox(s, x, y, cw, chh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(nd, GREEN_D, True)], size=15, align=PP_ALIGN.CENTER, first=True, space_after=0)
        if i < 2:
            _, ta = textbox(s, x + cw, y, gap, chh, anchor=MSO_ANCHOR.MIDDLE)
            para(ta, [("→", T3, True)], size=20, align=PP_ALIGN.CENTER, first=True, space_after=0)
        x += cw + gap
    _, ts = textbox(s, 0.9 + 3 * cw + 2 * gap + 0.1, y, 1.4, chh, anchor=MSO_ANCHOR.MIDDLE)
    para(ts, [("→ снапшот", T3, True)], size=13, first=True, space_after=0)
    bullets(s, 0.9, 3.9, 11.5, 2.8, [
        ("Декларативный сценарий: ", "фазы → шаги → чипсы. Контент отделён от движка."),
        ("Human-in-the-loop: ", "клиент отвечает чипсами, граф делает один ход и ждёт."),
        ("Stateless поверх Django ORM ", "— единый источник правды, без чекпойнтера LangGraph."),
        ("Состояние после рестарта ", "= phase + step_index + plan + budget + messages. Переживает перезапуск и refresh."),
    ], size=16, gap=14)


# 20 — AI/LLM
def s20():
    s = slide(); set_bg(s, PAPER); top_accent(s); footer(s, 20)
    kicker(s, "06 · Как это построено")
    title(s, "Язык тревоги, а не продажи"); title_rule(s)
    bullets(s, 0.9, 2.2, 6.6, 4.4, [
        ("claude-haiku-4-5 ", "генерирует каждое сообщение под контекст: погода, дети, отель."),
        ("Тон официальный, на «Вы»; ", "схема explain-then-offer (сначала зачем, потом предложение)."),
        ("На любой ошибке или без ключа ", "→ канонический fallback-текст шага."),
        ("Поэтому поездка всегда доходит до конца ", "— даже без сети и без API-ключа."),
    ], size=16.5, gap=16)
    # пример-цитата
    card(s, 7.9, 2.3, 4.5, 4.1, fill=CARD_BG, radius=0.05)
    rect(s, 7.9, 2.5, 0.10, 3.7, fill=GREEN)
    _, tq = textbox(s, 8.25, 2.65, 3.85, 3.5, anchor=MSO_ANCHOR.TOP)
    para(tq, [("Принцип системного промпта", GREEN_D, True)], size=13, first=True, space_after=10)
    para(tq, [("«Каждое сообщение отвечает на конкретную мысль-тревогу клиента — он чувствует заботу, не продажу. Стиль официальный, на «Вы». Без нагнетания.»", INK, False)],
         size=14.5, space_after=12, line=1.18)
    para(tq, [("→ fallback гарантирует, что демо никогда не «падает» из-за сети.", T2, False)],
         size=12.5, space_after=0, line=1.15)


# 21 — Спасибо
def s21():
    s = slide(); set_bg(s, HERO)
    rect(s, 0, 0, PW, 0.05, fill=GREEN)
    wordmark(s, 0.9, 0.78, scale=1.05)
    _, tf = textbox(s, 0.9, 2.7, 11.5, 2.6)
    para(tf, [("Halyk перестаёт ", INK), ("продавать билеты", T3)],
         size=40, first=True, space_after=8, line=1.1, font=SERIF)
    para(tf, [("— и начинает ", INK), ("вести клиента.", GREEN_D)],
         size=40, space_after=0, line=1.1, font=SERIF)
    _, tf2 = textbox(s, 0.9, 5.5, 11.5, 1.2)
    para(tf2, [("От покупки билета до возвращения домой — каждая тревога закрыта в нужный момент.", T2)],
         size=17, first=True, space_after=14)
    para(tf2, [("Спасибо!   ", INK, True), ("Демо: docker compose up → localhost:8000", T3, False)],
         size=15, space_after=0)


def main():
    for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
               s11, s12, s13, s14, s15, s16, s17, s18, s19, s20, s21]:
        fn()
    out = os.path.join(HERE, "Halyk_Smart_Travel.pptx")
    prs.save(out)
    print("OK saved:", out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
